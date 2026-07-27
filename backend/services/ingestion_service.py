"""Ingestion service — parse documents and store chunks + embeddings."""
from langchain_core.documents import Document
import io
import uuid
from typing import Any

import httpx
from sqlalchemy.ext.asyncio import AsyncSession

from core.config import settings
from models.db_models import DocumentChunk, KnowledgeAsset

from langchain_text_splitters import RecursiveCharacterTextSplitter
from services.ai_service import get_embeddings_batch

CHUNK_SIZE = 512
CHUNK_OVERLAP = 64

def _split_text(text: str) -> list[str]:
    """LangChain RecursiveCharacterTextSplitter."""
    text = text.replace("\x00", "")
    splitter = RecursiveCharacterTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)
    docs = splitter.create_documents([text])
    return [doc.page_content for doc in docs]

def _split_docs(docs: list[Document]) -> list[Document]:
    """LangChain RecursiveCharacterTextSplitter."""
    # Sanitize null bytes to prevent database write errors
    for doc in docs:
        if doc.page_content:
            doc.page_content = doc.page_content.replace("\x00", "")
            
    splitter = RecursiveCharacterTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)
    doc_split = splitter.split_documents(docs)
    return doc_split

async def _embed_texts(texts: list[str]) -> list[list[float]]:
    """Batch embed texts via unified AI service (OpenAI / Ollama)."""
    return await get_embeddings_batch(texts)


async def _store_chunks(
    asset: KnowledgeAsset,
    chunks: list[str],
    embeddings: list[list[float]],
    db: AsyncSession,
    extra_metadata: dict | None = None,
    chunk_metadatas: list[dict] | None = None,
) -> int:
    for i, (chunk_text, embedding) in enumerate(zip(chunks, embeddings)):
        meta = {**(extra_metadata or {}), "chunk_index": i}
        if chunk_metadatas and i < len(chunk_metadatas):
            meta.update(chunk_metadatas[i])
        chunk = DocumentChunk(
            asset_id=asset.id,
            chunk_text=chunk_text,
            chunk_index=i,
            embedding=embedding,
            metadata_=meta,
        )
        db.add(chunk)
    await db.commit()
    return len(chunks)


async def ingest_file(
    title: str,
    source_type: str,
    file_bytes: bytes,
    filename: str,
    uploaded_by: uuid.UUID,
    db: AsyncSession,
) -> dict:
    """Parse a PDF or image file and store chunks in pgvector."""
    raw_text = ""
    documentSplit = []  # list[Document] with per-chunk page_content + metadata

    if filename.lower().endswith(".pdf") or source_type == "pdf":
        print(f"[ingest] Processing PDF: {filename} ({len(file_bytes)} bytes)")

        # ── Stage 1: pypdf ────────────────────────────────────────────────────
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            is_encrypted = reader.is_encrypted
            print(f"[ingest][pypdf] pages={len(reader.pages)}, encrypted={is_encrypted}")
            page_docs: list[Document] = []
            for page_num, page in enumerate(reader.pages):
                page_text = (page.extract_text() or "").replace("\x00", "").strip()
                print(f"[ingest][pypdf] page {page_num+1}: {len(page_text)} chars")
                if page_text:
                    page_docs.append(
                        Document(
                            page_content=page_text,
                            metadata={"source": filename, "page": page_num + 1},
                        )
                    )
            if page_docs:
                documentSplit = _split_docs(page_docs)
                raw_text = "\n".join(doc.page_content for doc in page_docs)
                print(f"[ingest][pypdf] extracted {len(raw_text)} chars, {len(documentSplit)} chunks")
        except Exception as e:
            print(f"[ingest][pypdf] FAILED: {e}")

        # ── Stage 2: pdfplumber ───────────────────────────────────────────────
        if not raw_text.strip():
            try:
                import pdfplumber
                with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                    page_docs = []
                    for page_num, page in enumerate(pdf.pages):
                        page_text = (page.extract_text() or "").replace("\x00", "").strip()
                        print(f"[ingest][pdfplumber] page {page_num+1}: {len(page_text)} chars")
                        if page_text:
                            page_docs.append(
                                Document(
                                    page_content=page_text,
                                    metadata={"source": filename, "page": page_num + 1},
                                )
                            )
                if page_docs:
                    documentSplit = _split_docs(page_docs)
                    raw_text = "\n".join(doc.page_content for doc in page_docs)
                    print(f"[ingest][pdfplumber] extracted {len(raw_text)} chars, {len(documentSplit)} chunks")
                else:
                    print("[ingest][pdfplumber] no text found in any page")
            except Exception as e:
                print(f"[ingest][pdfplumber] FAILED: {e}")

        # ── Stage 3: PyMuPDF (fitz) — handles embedded fonts & complex PDFs ──
        if not raw_text.strip():
            try:
                import fitz  # PyMuPDF
                doc_fitz = fitz.open(stream=file_bytes, filetype="pdf")
                page_docs = []
                for page_num in range(len(doc_fitz)):
                    page = doc_fitz[page_num]
                    page_text = page.get_text("text").replace("\x00", "").strip()
                    print(f"[ingest][pymupdf] page {page_num+1}: {len(page_text)} chars")
                    if page_text:
                        page_docs.append(
                            Document(
                                page_content=page_text,
                                metadata={"source": filename, "page": page_num + 1},
                            )
                        )
                doc_fitz.close()
                if page_docs:
                    documentSplit = _split_docs(page_docs)
                    raw_text = "\n".join(doc.page_content for doc in page_docs)
                    print(f"[ingest][pymupdf] extracted {len(raw_text)} chars, {len(documentSplit)} chunks")
                else:
                    print("[ingest][pymupdf] no text found — PDF may be image-only/scanned")
            except Exception as e:
                print(f"[ingest][pymupdf] FAILED: {e}")

        # ── Stage 4: unstructured (OCR / fallback layout parser) ─────────────
        if not raw_text.strip():
            try:
                from unstructured.partition.auto import partition
                elements = partition(file=io.BytesIO(file_bytes), metadata_filename=filename)
                raw_text = "\n".join(str(e) for e in elements if str(e).strip())
                print(f"[ingest][unstructured] extracted {len(raw_text)} chars")
            except Exception as e:
                print(f"[ingest][unstructured] FAILED: {e}")

        if not raw_text.strip():
            print(f"[ingest] WARNING: All PDF extractors failed — PDF may be scanned/image-only with no OCR support")
            return {
                "asset_id": None,
                "chunks_stored": 0,
                "message": f"Could not extract text from {filename}. The PDF may be scanned or image-only.",
            }

    # ── PPTX / PPT files: python-pptx ─────────────────────────────────────────
    elif filename.lower().endswith(".pptx") or filename.lower().endswith(".ppt") or source_type == "pptx":
        print(f"[ingest] Processing PPTX: {filename} ({len(file_bytes)} bytes)")
        try:
            import pptx
            prs = pptx.Presentation(io.BytesIO(file_bytes))
            slide_docs: list[Document] = []
            for slide_num, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            t = (paragraph.text or "").replace("\x00", "").strip()
                            if t:
                                slide_texts.append(t)
                if slide_texts:
                    slide_docs.append(
                        Document(
                            page_content="\n".join(slide_texts),
                            metadata={"source": filename, "slide": slide_num + 1},
                        )
                    )
            if slide_docs:
                documentSplit = _split_docs(slide_docs)
                raw_text = "\n".join(doc.page_content for doc in slide_docs)
                print(f"[ingest][pptx] extracted {len(raw_text)} chars, {len(documentSplit)} chunks")
        except Exception as e:
            print(f"[ingest][pptx] FAILED: {e}")

    # ── DOCX / DOC files: python-docx ─────────────────────────────────────────
    elif filename.lower().endswith(".docx") or filename.lower().endswith(".doc") or source_type == "docx":
        print(f"[ingest] Processing DOCX: {filename} ({len(file_bytes)} bytes)")
        try:
            import docx
            doc = docx.Document(io.BytesIO(file_bytes))
            paragraphs = [(p.text or "").replace("\x00", "").strip() for p in doc.paragraphs if (p.text or "").strip()]
            if paragraphs:
                raw_text = "\n\n".join(paragraphs)
                print(f"[ingest][docx] extracted {len(raw_text)} chars")
        except Exception as e:
            print(f"[ingest][docx] FAILED: {e}")

    # ── Fallback for image / other files: unstructured partition ──────────────
    if not raw_text.strip() and not documentSplit:
        try:
            from unstructured.partition.auto import partition
            elements = partition(file=io.BytesIO(file_bytes), metadata_filename=filename)
            raw_text = "\n".join(str(e) for e in elements if str(e).strip())
            print(f"[ingest][unstructured fallback] extracted {len(raw_text)} chars")
        except Exception as e:
            print(f"[ingest][unstructured fallback] FAILED: {e}")

    if not raw_text.strip() and not documentSplit:
        print(f"[ingest] WARNING: Could not extract text from {filename}")
        return {
            "asset_id": None,
            "chunks_stored": 0,
            "message": f"Could not extract text from {filename}.",
        }

    asset = KnowledgeAsset(
        title=title,
        source_type=source_type,
        source_uri=filename,
        uploaded_by=uploaded_by,
    )
    db.add(asset)
    await db.flush()

    if documentSplit:
        chunks = [doc.page_content for doc in documentSplit]
        chunk_metadatas = [doc.metadata for doc in documentSplit]
    else:
        chunks = _split_text(raw_text)
        chunk_metadatas = None

    embeddings = await _embed_texts(chunks)
    count = await _store_chunks(asset, chunks, embeddings, db, {"filename": filename}, chunk_metadatas)

    return {"asset_id": asset.id, "chunks_stored": count, "message": f"Ingested {count} chunks from {filename}."}






def extract_youtube_id(url: str) -> str | None:
    import re
    match = re.search(r"(?:v=|\/embed\/|youtu\.be\/|\/v\/|\/e\/|watch\?v=|&v=)([^#\&\?]{11})", url)
    return match.group(1) if match else None


async def fetch_youtube_transcript(video_id: str) -> str:
    import asyncio
    from youtube_transcript_api import YouTubeTranscriptApi

    def _fetch():
        transcript_list = YouTubeTranscriptApi().list(video_id)
        try:
            transcript = transcript_list.find_transcript(["en", "en-GB", "en-US"])
        except Exception:
            try:
                first_transcript = next(iter(transcript_list))
                transcript = first_transcript.translate("en")
            except Exception as e:
                raise ValueError(f"No transcripts available to fetch or translate: {e}")

        lines = []
        for entry in transcript.fetch():
            start_sec = int(entry.start)
            minutes = start_sec // 60
            seconds = start_sec % 60
            timestamp = f"[{minutes:02d}:{seconds:02d}]"
            lines.append(f"{timestamp} {entry.text}")
        return "\n".join(lines)

    return await asyncio.to_thread(_fetch)


async def ingest_url(
    url: str,
    source_type: str,
    uploaded_by: uuid.UUID,
    db: AsyncSession,
) -> dict:
    """Fetch and ingest content from a URL or Wikipedia page."""
    documentSplit = []
    if "drive.google.com" in url and "/file/d/" in url:
        parts = url.split("/file/d/")
        if len(parts) > 1:
            doc_id = parts[1].split("/")[0]
            url = f"https://drive.google.com/uc?id={doc_id}&export=download"

    yt_id = extract_youtube_id(url)

    if source_type == "wiki":
        try:
            # pyrefly: ignore [missing-import]
            import wikipedia
            # Extract page title from URL or use URL as search term
            page_title = url.split("/wiki/")[-1].replace("_", " ") if "/wiki/" in url else url
            page = wikipedia.page(page_title)
            raw_text = page.content
            title = page.title
        except Exception as e:
            raise ValueError(f"Could not fetch Wikipedia page: {e}")
    elif source_type == "youtube" or yt_id:
        try:
            video_id = yt_id or url
            raw_text = await fetch_youtube_transcript(video_id)
            title = f"YouTube Video: {video_id}"
        except Exception as e:
            raise ValueError(f"Failed to fetch YouTube transcript: {e}")
    else:
        async with httpx.AsyncClient(timeout=30.0) as client:
            resp = await client.get(url, follow_redirects=True)
            resp.raise_for_status()

        content_type = resp.headers.get("content-type", "").lower()
        if "text/html" in content_type:
            try:
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(resp.text, "html.parser")
                raw_text = soup.get_text(separator="\n", strip=True)
                title = soup.title.string if soup.title else url
            except Exception:
                raw_text = resp.text
                title = url
        else:
            filename = url.split("/")[-1].split("?")[0] or "document"
            cd = resp.headers.get("content-disposition", "")
            if "filename=" in cd:
                import re
                match = re.search(r'filename="?([^";]+)"?', cd)
                if match:
                    filename = match.group(1)

            # Try OnlinePDFLoader from langchain_community if URL points to a PDF
            is_pdf = filename.lower().endswith(".pdf") or "application/pdf" in content_type
            if is_pdf:
                try:
                    from langchain_community.document_loaders import OnlinePDFLoader
                    loader = OnlinePDFLoader(url)
                    docs = loader.load()
                    documentSplit = _split_docs(docs)
                    raw_text = "\n".join(doc.page_content for doc in docs)
                    title = filename
                except Exception as e:
                    print(f"OnlinePDFLoader failed: {e}")

            if not raw_text.strip():
                try:
                    from unstructured.partition.auto import partition
                    elements = partition(file=io.BytesIO(resp.content), metadata_filename=filename)
                    raw_text = "\n".join(str(e) for e in elements if str(e).strip())
                    title = filename
                except Exception as e:
                    if content_type.startswith("text/"):
                        raw_text = resp.content.decode("utf-8", errors="ignore")
                        title = filename
                    else:
                        raise ValueError(f"Failed to parse binary document from URL: {e}")

    asset = KnowledgeAsset(
        title=title or url,
        source_type=source_type,
        source_uri=url,
        uploaded_by=uploaded_by,
    )
    db.add(asset)
    await db.flush()

    if documentSplit:
        chunks = [doc.page_content for doc in documentSplit]
        chunk_metadatas = [doc.metadata for doc in documentSplit]
    else:
        chunks = _split_text(raw_text)
        chunk_metadatas = None

    embeddings = await _embed_texts(chunks)
    count = await _store_chunks(asset, chunks, embeddings, db, {"source_url": url}, chunk_metadatas)

    return {"asset_id": asset.id, "chunks_stored": count, "message": f"Ingested {count} chunks from {url}."}
